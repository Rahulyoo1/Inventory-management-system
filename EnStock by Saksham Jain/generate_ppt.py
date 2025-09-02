from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_inventory_presentation(output_file):
    # Create presentation object
    prs = Presentation()

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Inventory Management System"
    subtitle.text = "A Python Django Project"

    # Customizing title and subtitle fonts and adding background color
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(230, 230, 250)  # Lavender background

    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)  # Blue text

    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.italic = True
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 0, 128)  # Purple text

    # Slide 2: Introduction Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Introduction"
    content.text = (
        "This project demonstrates an Inventory Management System built using Python and Django. "
        "It helps in managing products, stock levels, and transactions efficiently."
    )

    # Add background color
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(240, 248, 255)  # Alice blue

    # Slide 3: Features of the System
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Features"
    content.text = (
        "1. User Authentication and Authorization\n"
        "2. Add, Update, and Delete Inventory Items\n"
        "3. Real-Time Stock Tracking\n"
        "4. Reporting and Data Export\n"
        "5. Dashboard with Insights"
    )

    # Add background color and rectangle shapes
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(224, 255, 255)  # Light cyan
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5), Inches(9), Inches(0.4))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(135, 206, 235)  # Sky blue

    # Slide 4: Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title.text = "System Architecture"
    
    title_text_frame = title.text_frame
    title_text_frame.paragraphs[0].font.size = Pt(32)
    title_text_frame.paragraphs[0].font.bold = True

    # Adding shapes to represent the architecture
    db_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(2), Inches(2), Inches(4), Inches(1)
    )
    db_shape.fill.solid()
    db_shape.fill.fore_color.rgb = RGBColor(91, 155, 213)  # Blue
    db_shape.text = "Database"

    server_shape = slide.shapes.add_shape(
        MSO_SHAPE.CLOUD, Inches(7), Inches(2), Inches(2), Inches(1.5)
    )
    server_shape.fill.solid()
    server_shape.fill.fore_color.rgb = RGBColor(112, 173, 71)  # Green
    server_shape.text = "Server"

    # Slide 5: Technologies Used
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Technologies Used"
    content.text = (
        "- Python\n"
        "- Django Framework\n"
        "- SQLite Database\n"
        "- HTML, CSS, JavaScript (Frontend)\n"
        "- Bootstrap for Responsive Design"
    )

    # Add gradient background
    slide.background.fill.gradient()
    gradient_stops = slide.background.fill.gradient_stops
    gradient_stops[0].position = 0.0
    gradient_stops[0].color.rgb = RGBColor(255, 228, 196)  # Bisque
    gradient_stops[1].position = 1.0
    gradient_stops[1].color.rgb = RGBColor(255, 140, 0)  # Dark orange

    # Slide 6: Demo Screenshots
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Content layout without images
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Demo Screenshots"
    content.text = "Screenshots can be added here to showcase the project UI and functionality."

    # Add dashed line separator
    line = slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, Inches(0.5), Inches(1.5), Inches(8.5), Inches(0))
    line.line.color.rgb = RGBColor(255, 69, 0)  # Red orange
    line.line.dash_style = 3  # Dashed line

    # Slide 7: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Conclusion"
    content.text = (
        "The Inventory Management System built using Django simplifies the tracking and management of inventory.\n"
        "It is scalable, efficient, and user-friendly."
    )

    # Add subtle gray background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(211, 211, 211)  # Light gray

    # Save the presentation
    prs.save(output_file)

# Generate the presentation
output_file = "Inventory_Management_Project.pptx"
create_inventory_presentation(output_file)
print(f"Presentation saved as {output_file}")
